"""
PowerPoint export utilities.

Benjamin Mallada 2025
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pathlib import Path
from PIL import Image


def export_images_to_pptx(image_directory, output_presentation_path, images_per_slide=4):
    """Export BMP images from ``image_directory`` into a PowerPoint file."""
    image_directory = Path(image_directory)
    bmp_files = sorted(image_directory.glob("*.bmp"))

    if not bmp_files:
        print("No BMP files found.")
        return

    presentation = Presentation()
    slide_width = presentation.slide_width = Inches(10)
    slide_height = presentation.slide_height = Inches(10)

    # Determine layout grid based on the number of images per slide
    layout_grid_map = {
        1: (1, 1),
        2: (1, 2),
        4: (2, 2),
        6: (2, 3),
        9: (3, 3),
    }
    rows, columns = layout_grid_map.get(images_per_slide, (2, 2))
    margin = Inches(0.3)
    row_spacing = Inches(0.1)
    # Space allocated for each image placeholder
    image_width = (slide_width - (columns + 1) * margin) / columns
    image_height = (slide_height - (rows + 1) * margin) / rows

    slide = None
    for index, bmp_file in enumerate(bmp_files):
        if index % images_per_slide == 0:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        row_index = (index % images_per_slide) // columns
        col_index = (index % images_per_slide) % columns

        left_position = margin + col_index * (image_width + margin)
        top_position = margin + row_index * (image_height + margin + row_spacing)
        # Insert image
        try:
            slide.shapes.add_picture(
                str(bmp_file),
                left_position,
                top_position,
                width=image_width,
                height=image_height,
            )
            # Adjust vertical position to center image if it becomes shorter

        except Exception as error:  # pragma: no cover - printing error
            print(f"Could not insert {bmp_file.name}: {error}")

    presentation.save(output_presentation_path)
    print(f"✔ Saved PowerPoint: {output_presentation_path}")

